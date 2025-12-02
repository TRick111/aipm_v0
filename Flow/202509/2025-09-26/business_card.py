from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# ロゴ（円で代替）
shape = slide.shapes.add_shape(
    1,  # msoShapeOval
    Inches(6.6), Inches(0.25), Inches(0.95), Inches(0.95)
)
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 112, 192)

# 右上の法人名
textbox = slide.shapes.add_textbox(Inches(5.1), Inches(0.3), Inches(3), Inches(1.5))
tf = textbox.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "独立行政法人 国立病院機構\n横浜医療センター"
p.font.size = Pt(11)
p.font.name = "MS Mincho"
p.alignment = PP_ALIGN.LEFT

# 薬剤師肩書き
textbox = slide.shapes.add_textbox(Inches(1), Inches(2.1), Inches(2), Inches(0.5))
tf = textbox.text_frame
p = tf.add_paragraph()
p.text = "薬剤師"
p.font.size = Pt(11)
p.font.name = "MS Mincho"

# 名前
textbox = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(3), Inches(1))
tf = textbox.text_frame
p = tf.add_paragraph()
p.text = "関口 華乃"
p.font.size = Pt(21)
p.font.name = "MS Mincho"

p = tf.add_paragraph()
p.text = "せきぐち かの"
p.font.size = Pt(9)
p.font.name = "MS Gothic"

# 住所・連絡先
textbox = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(5.8), Inches(2))
tf = textbox.text_frame
p = tf.add_paragraph()
p.text = "〒245-8575 神奈川県横浜市戸塚区原宿3-60-2"
p.font.size = Pt(9)
p.font.name = "MS Gothic"

p = tf.add_paragraph()
p.text = "TEL: 045-851-2621（代表）"
p.font.size = Pt(9)
p.font.name = "MS Gothic"

p = tf.add_paragraph()
p.text = "E-mail: sekiguchi.kano.rb@mail.hosp.go.jp"
p.font.size = Pt(9)
p.font.name = "MS Gothic"

prs.save('名刺再現_最終調整版.pptx')
